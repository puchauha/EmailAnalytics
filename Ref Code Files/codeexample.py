"""
POC: LLM-driven parsing of outage emails → Vector DB → Executive Summary → PPT report
"""

import os
import re
import json
import uuid
from pathlib import Path
from collections import Counter
from tqdm import tqdm

# Vector DB & Embedding
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

# LLM + Visualization + PPT
import openai
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# -------------------------------
# CONFIG
# -------------------------------
DATA_DIR = Path("./outage_texts")
CHROMA_DIR = "./chroma_llm"
COLLECTION_NAME = "partner_outages"
PPTX_FILE = "Outage_Summary_POC.pptx"
CHART_DIR = Path("./charts_llm")
CHART_DIR.mkdir(exist_ok=True, parents=True)

openai.api_key = os.getenv("OPENAI_API_KEY")
EMBED_MODEL = "all-MiniLM-L6-v2"

# -------------------------------
# Initialize
# -------------------------------
embedder = SentenceTransformer(EMBED_MODEL)
client = chromadb.Client(Settings(chroma_db_impl="duckdb+parquet", persist_directory=CHROMA_DIR))

if COLLECTION_NAME in [c.name for c in client.list_collections()]:
    collection = client.get_collection(COLLECTION_NAME)
else:
    collection = client.create_collection(COLLECTION_NAME)

# -------------------------------
# LLM Helper Functions
# -------------------------------
def llm_parse_email(email_text: str):
    """
    Use LLM to extract partner name, date, reason from outage mail.
    Returns a dict.
    """
    prompt = f"""
You are a system that extracts structured data from outage emails.
Given the text below, extract the following fields:
1. partner_name
2. outage_date
3. outage_reason (briefly summarize, e.g., "VPN failure", "Database outage")
4. impact_summary (1-2 lines summarizing what happened)

Return a JSON object with keys: partner_name, outage_date, outage_reason, impact_summary.

Email:
---
{email_text}
---
If unsure, return "unknown" for the missing fields.
"""
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role":"system","content":"You extract structured outage data from text emails."},
                {"role":"user","content":prompt}
            ],
            temperature=0
        )
        text = response["choices"][0]["message"]["content"].strip()
        data = json.loads(text)
        return data
    except Exception as e:
        print("Parsing failed:", e)
        return {"partner_name": "unknown", "outage_date": "unknown", "outage_reason": "unknown", "impact_summary": ""}


def llm_generate_summary(partner: str, all_texts: list):
    """
    Generate executive summary for partner outages
    """
    joined = "\n---\n".join(all_texts[:6])
    prompt = f"""
You are creating an executive summary of outage notifications from partner '{partner}'.

Summarize in 4-5 sentences:
- Main outage reasons and trends
- Notable impacts or patterns
- Any recommendations or next steps

Keep it formal and concise.
Messages:
{joined}
"""
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[{"role":"system","content":"You summarize outage communications professionally."},
                      {"role":"user","content":prompt}],
            temperature=0.3,
            max_tokens=300
        )
        return response["choices"][0]["message"]["content"].strip()
    except Exception as e:
        print("Summary generation failed:", e)
        return "Summary unavailable due to API error."

# -------------------------------
# Step 1: Ingest emails
# -------------------------------
def ingest_emails():
    txt_files = sorted(DATA_DIR.glob("*.txt"))
    if not txt_files:
        print("No emails found in", DATA_DIR)
        return

    for f in tqdm(txt_files, desc="Parsing and embedding emails"):
        text = f.read_text(encoding="utf8", errors="ignore")
        parsed = llm_parse_email(text)
        emb = embedder.encode(text).tolist()

        metadata = {
            "filename": f.name,
            "partner": parsed.get("partner_name", "unknown"),
            "date": parsed.get("outage_date", "unknown"),
            "reason": parsed.get("outage_reason", "unknown"),
            "impact": parsed.get("impact_summary", "")
        }

        collection.add(
            ids=[str(uuid.uuid4())],
            documents=[text],
            metadatas=[metadata],
            embeddings=[emb]
        )
    client.persist()
    print("✅ Ingest complete.")

# -------------------------------
# Step 2: Generate PPT per partner
# -------------------------------
def generate_ppt():
    all_meta = collection.get(include=["metadatas"])
    flat_meta = []
    for m in all_meta.get("metadatas", []):
        if isinstance(m, list):
            flat_meta.extend(m)
        else:
            flat_meta.append(m)

    partners = sorted({m.get("partner", "unknown") for m in flat_meta})
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Partner Outage Summary"
    cover.placeholders[1].text = "Generated automatically using LLM + Vector DB"

    for partner in tqdm(partners, desc="Creating slides"):
        results = collection.query(where={"partner": partner}, n_results=100, include=["documents", "metadatas"])
        docs, metas = results["documents"], results["metadatas"]

        # Flatten
        if docs and isinstance(docs[0], list):
            docs = [d for sub in docs for d in sub]
            metas = [m for sub in metas for m in sub]

        reasons = [m.get("reason","unknown") for m in metas]
        counts = Counter(reasons)

        # Chart
        plt.figure(figsize=(6,3))
        plt.bar(counts.keys(), counts.values())
        plt.title(f"Outage Reasons — {partner}")
        plt.ylabel("Count")
        plt.tight_layout()
        chart_path = CHART_DIR / f"{partner}_chart.png"
        plt.savefig(chart_path, bbox_inches="tight")
        plt.close()

        # Summary
        summary = llm_generate_summary(partner, docs)

        # Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title.text = f"{partner} — Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True

        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
        body_tf = body.text_frame
        body_tf.text = summary
        body_tf.word_wrap = True
        body_tf.paragraphs[0].font.size = Pt(12)

        slide.shapes.add_picture(str(chart_path), Inches(6.3), Inches(1.2), width=Inches(3.8), height=Inches(3.0))

    prs.save(PPTX_FILE)
    print("✅ PPT generated:", PPTX_FILE)

# -------------------------------
# Run all
# -------------------------------
if __name__ == "__main__":
    ingest_emails()
    generate_ppt()
