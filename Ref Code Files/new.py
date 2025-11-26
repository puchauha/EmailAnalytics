# %%
"""
POC: LLM-driven parsing of outage emails → JSON → PPT report
"""
import os
import json
import re
from pathlib import Path
from collections import Counter
from tqdm import tqdm
from dotenv import load_dotenv
from openai import OpenAI
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# %%
load_dotenv()
DATA_DIR = Path(os.getenv("DATA_DIR", "./logistics_outage_emails"))
CHART_DIR = Path(os.getenv("CHART_DIR", "./charts"))
OUTPUT_JSON = Path(os.getenv("OUTPUT_JSON", "./parsed_outages.json"))
PPTX_FILE = Path(os.getenv("PPTX_FILE", "./Partner_Outage_Summary.pptx"))

CHART_DIR.mkdir(exist_ok=True, parents=True)

client = OpenAI()

# %%
def llm_parse_email(email_text: str):
    """Use LLM to extract structured data from outage mail."""
    prompt = f"""
You are a system that extracts structured data from outage emails.
Given the text below, extract the following fields:
1. partner_name
2. outage_date
3. outage_reason (briefly summarize, e.g., "VPN failure", "Database outage")
4. impact_summary (1-2 lines summarizing what happened)

Return a JSON object only — no extra text or markdown.

Email:
---
{email_text}
---
If unsure, return "unknown" for missing fields.
"""
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You extract structured outage data from text emails."},
                {"role": "user", "content": prompt}
            ],
            temperature=0
        )
        text = response.choices[0].message.content.strip()

        # Extract JSON safely
        match = re.search(r"\{.*\}", text, re.DOTALL)
        if match:
            text = match.group(0)

        data = json.loads(text)
        return data
    except Exception as e:
        print("⚠️ Parsing failed:", e)
        return {
            "partner_name": "unknown",
            "outage_date": "unknown",
            "outage_reason": "unknown",
            "impact_summary": ""
        }

# %%
def llm_generate_summary(partner: str, all_texts: list):
    """Generate executive summary for partner outages."""
    joined = "\n---\n".join(all_texts[:6])
    prompt = f"""
You are creating an executive summary of outage notifications from partner '{partner}'.

Summarize in 4-5 sentences:
- Main outage reasons and trends
- Notable impacts or patterns
- Any recommendations or next steps

Keep it formal and concise.
Messages:
{joined}
"""
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You summarize outage communications professionally."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=300
        )
        summary = response.choices[0].message.content.strip()
        print(f"✅ Summary generated for {partner}")
        return summary
    except Exception as e:
        print("⚠️ Summary generation failed:", e)
        return "Summary unavailable due to API error."

# %%
def ingest_emails():
    """Parse emails and store structured data in JSON."""
    print("Ingesting emails from", DATA_DIR)
    txt_files = sorted(DATA_DIR.glob("*.txt"))
    if not txt_files:
        print("No emails found in", DATA_DIR)
        return []

    parsed_records = []

    for f in tqdm(txt_files, desc="Parsing outage emails"):
        text = f.read_text(encoding="utf8", errors="ignore")
        parsed = llm_parse_email(text)
        parsed["filename"] = f.name
        parsed["raw_text"] = text
        parsed_records.append(parsed)

    # Save to JSON
    with open(OUTPUT_JSON, "w", encoding="utf8") as fp:
        json.dump(parsed_records, fp, indent=2, ensure_ascii=False)

    print(f"✅ Parsed data saved to {OUTPUT_JSON}")
    return parsed_records

# %%
def generate_ppt(records):
    """Create PPT using parsed JSON records."""
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Partner Outage Summary"
    cover.placeholders[1].text = "Generated automatically using LLM (JSON mode)"

    # Group by partner
    partners = {}
    for r in records:
        partners.setdefault(r["partner_name"], []).append(r)

    for partner, recs in tqdm(partners.items(), desc="Creating slides"):
        reasons = [r["outage_reason"] for r in recs]
        counts = Counter(reasons)

        # Chart
        plt.figure(figsize=(6, 3))
        plt.bar(counts.keys(), counts.values())
        plt.title(f"Outage Reasons — {partner}")
        plt.ylabel("Count")
        plt.tight_layout()
        chart_path = CHART_DIR / f"{partner}_chart.png"
        plt.savefig(chart_path, bbox_inches="tight")
        plt.close()

        # Summary
        texts = [r["raw_text"] for r in recs]
        summary = llm_generate_summary(partner, texts)

        # Slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title.text = f"{partner} — Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        title.text_frame.paragraphs[0].font.bold = True

        # Body text
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
        body_tf = body.text_frame
        body_tf.text = summary
        body_tf.word_wrap = True
        body_tf.paragraphs[0].font.size = Pt(12)

        # Chart image
        slide.shapes.add_picture(str(chart_path), Inches(6.3), Inches(1.2),
                                 width=Inches(3.8), height=Inches(3.0))

    prs.save(PPTX_FILE)
    print("✅ PPT generated:", PPTX_FILE)

# %%
if __name__ == "__main__":
    records = ingest_emails()
    if records:
        generate_ppt(records)
